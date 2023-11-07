import io
import traceback
from typing import IO
from pptx import Presentation
from pptx.util import Inches,Pt
import pytesseract
from PIL import Image
import ai.get_response
import aspose.slides as slides
import aspose.pydrawing as drawing


class GetText:
    def __init__(self,file:IO):
        self.file = file


    async def pptx_method(self):
        text = ""
        presentation = Presentation(self.file)
        print("pptxmethod内部")
        try:
            for slide in presentation.slides:
                for shape in slide.shapes:
                    if hasattr(shape, 'text'):
                        text += shape.text + '\n'
                    if shape.shape_type == 13:
                        image = shape.image
                        try:
                            img = Image.open(io.BytesIO(image.blob))
                            img = img.convert('L')
                            img_text = pytesseract.image_to_string(img)
                            text += img_text + '\n'
                        except Exception as e:
                            print("有图片异常")
            # 检查是否有text
            if not text:
                print("no text")
            
            self.text = text
            #存数据库
            print("pptx发送ai")
            result = await self.send_to_ai(self.text)
            return result            
        
        except Exception as e:
            print("读取文件内容异常")
            print(traceback.format_exc())
            return "读取文件异常"

    async def ppt_method(self):
        print("ppt file")
        return "无法解析ppt文件，正在加紧开发，请尝试使用其他网站转换文件格式或者使用PowerPoint将ppt文件改成pptx文件再使用"
    
    async def send_to_ai(self,alltext):
        print("send to ai")
        operate =ai.get_response.response(alltext)
        result = await operate.complete()
        return result
        


